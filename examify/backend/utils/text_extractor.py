"""
Examify Text Extraction Utilities
Handles text extraction from PDF, PPT, and PPTX files.
"""

import re
import logging
from typing import List, Dict

logger = logging.getLogger(__name__)


def extract_text_from_pdf(filepath: str) -> str:
    """
    Extract text from a PDF file using PyPDF2.
    Returns cleaned, structured text.
    """
    try:
        import PyPDF2
        
        text_parts = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            total_pages = len(reader.pages)
            logger.info(f"Extracting text from PDF with {total_pages} pages")
            
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        cleaned = clean_text(page_text)
                        if len(cleaned) > 50:  # Skip nearly empty pages
                            text_parts.append(f"[Section {page_num + 1}]\n{cleaned}")
                except Exception as e:
                    logger.warning(f"Could not extract page {page_num + 1}: {e}")
                    continue
        
        return '\n\n'.join(text_parts)
    
    except ImportError:
        raise RuntimeError("PyPDF2 not installed. Run: pip install PyPDF2")
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        raise RuntimeError(f"Failed to extract text from PDF: {str(e)}")


def extract_text_from_pptx(filepath: str) -> str:
    """
    Extract text from a PowerPoint file using python-pptx.
    Returns cleaned, structured text per slide.
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(filepath)
        text_parts = []
        slide_count = len(prs.slides)
        logger.info(f"Extracting text from PPTX with {slide_count} slides")
        
        for slide_num, slide in enumerate(prs.slides):
            slide_texts = []
            
            # Extract slide title
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = clean_text(slide.shapes.title.text)
                if title:
                    slide_texts.append(f"Title: {title}")
            
            # Extract all text from shapes
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                shape_text = []
                for para in shape.text_frame.paragraphs:
                    para_text = ' '.join(run.text for run in para.runs if run.text.strip())
                    if para_text.strip():
                        cleaned = clean_text(para_text)
                        if cleaned:
                            shape_text.append(cleaned)
                
                if shape_text:
                    slide_texts.extend(shape_text)
            
            # Extract text from tables
            for shape in slide.shapes:
                if shape.shape_type == 19:  # Table
                    for row in shape.table.rows:
                        row_text = ' | '.join(
                            clean_text(cell.text) for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)
            
            if slide_texts:
                combined = '\n'.join(slide_texts)
                if len(combined) > 50:
                    text_parts.append(f"[Slide {slide_num + 1}]\n{combined}")
        
        return '\n\n'.join(text_parts)
    
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        raise RuntimeError(f"Failed to extract text from PPTX: {str(e)}")


def clean_text(text: str) -> str:
    """
    Clean and normalize extracted text.
    Removes page numbers, symbols, redundant whitespace, etc.
    """
    if not text:
        return ''
    
    # Remove page numbers (standalone numbers)
    text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)
    
    # Remove common slide artifacts
    text = re.sub(r'^\s*slide\s*\d+\s*$', '', text, flags=re.MULTILINE | re.IGNORECASE)
    text = re.sub(r'^\s*page\s*\d+\s*(of\s*\d+)?\s*$', '', text, flags=re.MULTILINE | re.IGNORECASE)
    
    # Remove URLs
    text = re.sub(r'http[s]?://\S+', '', text)
    
    # Remove excessive special characters but keep punctuation
    text = re.sub(r'[^\w\s.,;:!?\'\"()\-–—\[\]{}%$@#&*+=/<>]', ' ', text)
    
    # Remove lines that are just symbols or very short (likely headers/footers)
    lines = text.split('\n')
    cleaned_lines = []
    for line in lines:
        stripped = line.strip()
        # Keep lines that have at least 3 words or meaningful content
        word_count = len(stripped.split())
        if word_count >= 3 or (word_count > 0 and len(stripped) > 20):
            cleaned_lines.append(stripped)
    
    text = '\n'.join(cleaned_lines)
    
    # Normalize whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    
    return text.strip()


def split_into_chunks(text: str, max_chunk_size: int = 3000) -> List[str]:
    """
    Split large text into manageable chunks for AI processing.
    Tries to split at section/paragraph boundaries.
    """
    if len(text) <= max_chunk_size:
        return [text]
    
    chunks = []
    
    # First try splitting by section markers
    sections = re.split(r'\[(?:Section|Slide)\s+\d+\]', text)
    sections = [s.strip() for s in sections if s.strip() and len(s.strip()) > 100]
    
    current_chunk = []
    current_size = 0
    
    for section in sections:
        if current_size + len(section) > max_chunk_size and current_chunk:
            chunks.append('\n\n'.join(current_chunk))
            current_chunk = [section]
            current_size = len(section)
        else:
            current_chunk.append(section)
            current_size += len(section)
    
    if current_chunk:
        chunks.append('\n\n'.join(current_chunk))
    
    # If sections are too large, split by paragraphs
    final_chunks = []
    for chunk in chunks:
        if len(chunk) <= max_chunk_size:
            final_chunks.append(chunk)
        else:
            # Split by paragraphs
            paragraphs = chunk.split('\n\n')
            sub_chunk = []
            sub_size = 0
            for para in paragraphs:
                if sub_size + len(para) > max_chunk_size and sub_chunk:
                    final_chunks.append('\n\n'.join(sub_chunk))
                    sub_chunk = [para]
                    sub_size = len(para)
                else:
                    sub_chunk.append(para)
                    sub_size += len(para)
            if sub_chunk:
                final_chunks.append('\n\n'.join(sub_chunk))
    
    logger.info(f"Split text into {len(final_chunks)} chunks")
    return final_chunks if final_chunks else [text[:max_chunk_size]]


def extract_text(filepath: str, file_type: str) -> str:
    """
    Main extraction dispatcher based on file type.
    """
    file_type = file_type.lower().strip('.')
    
    if file_type == 'pdf':
        return extract_text_from_pdf(filepath)
    elif file_type in ('ppt', 'pptx'):
        return extract_text_from_pptx(filepath)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
